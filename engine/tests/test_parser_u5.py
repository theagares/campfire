"""
tests/test_parser_u5.py
U5(파서 확장: HWP/HWPX/XLSX/PPTX) 통합 테스트 (PLAN §6, §9.2, §10 Phase 5).

- XLSX/PPTX: openpyxl/python-pptx 로 실제 샘플을 생성해 parse_document + 전체
  파이프라인(run_pipeline)까지 실 PII 탐지/마스킹이 되는지 검증.
- HWPX: SECUREDOC_TEST_HWPX=1 로 켰고 한/글(한컴오피스)이 설치돼 있으면 pyhwpx 로
  실제 hwpx 샘플을 생성해 왕복 검증, 아니면 "미지원 통과" 경로만 검증.
  기본이 꺼짐인 이유는 아래 HWPX_OPT_IN 주석 참고(스위트가 멎었던 실측 사례).
- HWP: LibreOffice 미설치 환경에서 "우아한 미지원 통과"(PLAN §9.2, §11)를 검증.
  (LibreOffice 가 설치된 환경이면 실제 변환 경로까지 확인.)
"""

from __future__ import annotations

import asyncio
import io
import os
import subprocess
import sys
import tempfile

import pytest

from app.core import model_status
from app.core.detectors import registry
from app.core.parser import STATUS_FAILED, STATUS_OK, STATUS_UNSUPPORTED, parse_document
from app.core.parser import hwp as hwp_module
from app.core.pipeline.orchestrator import run_pipeline

_needs_models = pytest.mark.skipif(
    not model_status.all_ready(), reason="PII/인젝션 모델이 로컬에 없어 마스킹 결과를 검증할 수 없음"
)

XLSX_MIME = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


# ── XLSX ──────────────────────────────────────────────────────────────────────


def _make_xlsx_bytes() -> bytes:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "고객목록"
    ws.append(["구분", "메모"])
    # PERSON_NAME 룰(app/rules/pii_ko.yaml)은 "성명/이름/고객명 : 값" 앵커가
    # 같은 텍스트 런 안에 있어야 매치되므로, 헤더/값을 별도 셀로 나누지 않고
    # 한 셀 안에 라벨+값을 함께 넣는다(실제 스프레드시트의 "메모" 컬럼과 유사).
    ws.append(["신규 고객", "성명: 홍길동, 이메일: hong@example.com, 연락처: 010-1234-5678"])
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def test_xlsx_parse_extracts_all_sheet_text():
    text, status, reason = parse_document(_make_xlsx_bytes(), XLSX_MIME, "customers.xlsx")
    assert status == STATUS_OK
    assert reason is None
    assert "홍길동" in text
    assert "hong@example.com" in text
    assert "010-1234-5678" in text


@_needs_models
def test_xlsx_pipeline_detects_and_masks_pii():
    # encoder/llm_mcp 는 실제 subprocess + asyncio 파이프를 쓴다 — 이전 테스트의
    # (이미 닫힌) asyncio.run() 이벤트 루프에 묶인 캐시된 detector 를 재사용하면
    # 깨지므로(실측), 이 파일의 각 ML 테스트는 먼저 캐시를 비운다.
    registry.reset_cache()
    result = asyncio.run(
        run_pipeline(file_bytes=_make_xlsx_bytes(), mime_type=XLSX_MIME, file_name="customers.xlsx")
    )
    assert result["scanStatus"] == STATUS_OK
    # 실 PII 인코더가 이름/이메일/연락처 중 정확히 몇 개를 잡는지는 문맥(셀 배치,
    # 주변 텍스트)에 따라 실측으로 갈렸다(모델 정확도 이슈, 룰베이스처럼 결정적이지
    # 않음) — 여기서는 "파이프라인이 실 모델을 거쳐 최소 1건 이상 마스킹한다"만
    # 확인한다. 특정 필드 탐지를 강제하지 않는다.
    assert result["stats"]["piiCount"] >= 1
    assert result["maskedText"] != result["originalText"]


# ── PPTX ──────────────────────────────────────────────────────────────────────


def _make_pptx_bytes() -> bytes:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "고객 정보"
    body = slide.placeholders[1]
    body.text_frame.text = "성명: 김철수"
    p = body.text_frame.add_paragraph()
    p.text = "연락처: 010-9876-5432 / 이메일: kim@example.com"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_pptx_parse_extracts_all_slide_text():
    text, status, reason = parse_document(_make_pptx_bytes(), PPTX_MIME, "briefing.pptx")
    assert status == STATUS_OK
    assert reason is None
    assert "김철수" in text
    assert "010-9876-5432" in text
    assert "kim@example.com" in text


@_needs_models
def test_pptx_pipeline_detects_and_masks_pii():
    registry.reset_cache()
    result = asyncio.run(
        run_pipeline(file_bytes=_make_pptx_bytes(), mime_type=PPTX_MIME, file_name="briefing.pptx")
    )
    assert result["scanStatus"] == STATUS_OK
    # 실 PII 인코더가 이름/이메일/연락처 중 정확히 몇 개를 잡는지는 문맥에 따라
    # 실측으로 갈렸다(모델 정확도 이슈) — "최소 1건 이상 마스킹"만 확인한다.
    assert result["stats"]["piiCount"] >= 1
    assert result["maskedText"] != result["originalText"]


# ── HWP (LibreOffice 변환) ────────────────────────────────────────────────────


def test_hwp_graceful_unsupported_or_real_conversion():
    """LibreOffice 미설치 환경 → unsupported 로 우아하게 통과(PLAN §9.2/§11).

    LibreOffice 가 설치돼 있으면(예외적) 더미 바이트는 진짜 HWP가 아니므로
    변환 실패 → failed 로 통과하는 것도 정상(§9.2 "파싱 실패"도 동일 정책).
    """
    result = asyncio.run(
        run_pipeline(
            file_bytes=b"this is not a real hwp binary",
            mime_type="application/x-hwp",
            file_name="report.hwp",
        )
    )

    if hwp_module._find_soffice() is None:
        assert result["scanStatus"] == STATUS_UNSUPPORTED
        assert "LibreOffice" in (result["reason"] or "")
    else:
        assert result["scanStatus"] in (STATUS_UNSUPPORTED, STATUS_FAILED)

    assert result["stats"]["piiCount"] == 0  # 검사 자체를 안 함(§9.2)


def test_hwp_find_soffice_returns_none_or_path():
    # 탐색 함수 자체가 예외 없이 동작하는지만 확인(있으면 str, 없으면 None)
    found = hwp_module._find_soffice()
    assert found is None or isinstance(found, str)


# ── HWPX (pyhwpx) ─────────────────────────────────────────────────────────────


# 한/글 자동화는 **옵트인**이다. 기본으로 켜두면 안 되는 이유가 실측으로 확인됐다:
#
#   Hwp(new=True, visible=False, register_module=True) 가 한/글의 보안 모듈 등록 승인
#   대화상자를 띄우고 거기서 무한 대기한다. visible=False 라 사용자에게는 그 창이
#   보이지도 않고, pyhwpx 에는 타임아웃이 없다 — pytest 가 영영 끝나지 않는다
#   (실측: Hwp.exe 가 10분 넘게 응답 없음, 전체 스위트가 그 지점에서 멈춤).
#
# CI 에서는 pyhwpx 가 없어 ImportError 로 조용히 skip 되기 때문에 이 문제가 오래
# 드러나지 않았다. 하필 한/글이 깔린 개발자 기계에서만 터진다.
#
# 그래서 SECUREDOC_TEST_HWPX=1 일 때만 시도한다. 켠 경우에도 아래 _spawn_hwpx_sample()
# 이 별도 프로세스 + 타임아웃으로 감싸므로, 대화상자가 떠도 스위트가 멎지는 않는다.
HWPX_OPT_IN = os.environ.get("SECUREDOC_TEST_HWPX") == "1"
HWPX_TIMEOUT_SEC = float(os.environ.get("SECUREDOC_TEST_HWPX_TIMEOUT_SEC", "120"))


# 자식 프로세스에서 돌릴 샘플 생성 스크립트. 이 코드를 이 프로세스에서 직접 부르면
# 승인 대화상자에 걸렸을 때 pytest 를 통째로 붙잡는다(위 주석 참고) — 프로세스를
# 분리해야 타임아웃으로 끊을 수 있다.
_HWPX_CHILD = r"""
import os, sys
out = sys.argv[1]
try:
    from pyhwpx import Hwp
except ImportError:
    sys.exit(3)
hwp = None
try:
    hwp = Hwp(new=True, visible=False, register_module=True)
    hwp.insert_text("성명: 박영희")
    hwp.insert_text(" 이메일: park@example.com")
    if not hwp.save_as(out, format="HWPX") or not os.path.exists(out):
        sys.exit(4)
except Exception:
    sys.exit(5)
finally:
    try:
        if hwp is not None:
            hwp.quit(save=False)
    except Exception:
        pass
sys.exit(0)
"""


def _try_make_real_hwpx_bytes() -> bytes | None:
    """실제 hwpx 샘플을 만들어 반환. 만들 수 없으면 None(미지원 통과 경로만 검증).

    한/글 자동화를 **자식 프로세스**에서 돌리고 타임아웃을 건다. 승인 대화상자에 걸려도
    이 함수는 시간 내에 돌아오고, 그 경우 테스트는 skip 된다.
    """
    if not HWPX_OPT_IN:
        return None

    fd, path = tempfile.mkstemp(suffix=".hwpx")
    os.close(fd)
    os.remove(path)

    try:
        proc = subprocess.run(
            [sys.executable, "-c", _HWPX_CHILD, path],
            timeout=HWPX_TIMEOUT_SEC,
            capture_output=True,
        )
    except subprocess.TimeoutExpired:
        # 자식은 죽였지만 COM 으로 뜬 Hwp.exe 는 남을 수 있다. 사용자가 열어둔 한/글
        # 창까지 죽일 수는 없으므로 여기서 정리하지 않고 알리기만 한다.
        pytest.skip(
            f"한/글 자동화가 {HWPX_TIMEOUT_SEC:.0f}초 안에 끝나지 않았습니다 "
            "(보안 모듈 등록 승인 대화상자 가능성). 남아 있는 Hwp.exe 를 확인하세요."
        )

    if proc.returncode != 0 or not os.path.exists(path):
        if os.path.exists(path):
            os.remove(path)
        return None

    try:
        with open(path, "rb") as f:
            return f.read()
    finally:
        if os.path.exists(path):
            os.remove(path)


def test_hwpx_real_extraction_or_graceful_unsupported():
    hwpx_bytes = _try_make_real_hwpx_bytes()

    if hwpx_bytes is None:
        # 실제 샘플을 못 만든 경우(옵트인 꺼짐/한/글 미설치/생성 실패) — 여기서 지킬 건
        # "§9.2: 파서는 예외를 밖으로 던지지 않는다" 하나다.
        #
        # 예전엔 이 분기가 status ∈ {unsupported, failed} 를 단언했는데, 그건 **한/글이
        # 없는 기계**를 전제한 것이었다. 한/글이 깔려 있으면 COM 이 이 쓰레기 바이트를
        # 평문 문서로 열어버려서 status=ok, text='not a real hwpx' 가 나온다(실측).
        # 그 전제가 깨진 걸 아무도 못 봤던 이유는, 그때까지 이 테스트가 위쪽
        # _try_make_real_hwpx_bytes() 의 승인 대화상자에서 멈춰 여기까지 오지도 못했기
        # 때문이다(HWPX_OPT_IN 주석 참고).
        text, status, reason = parse_document(
            b"not a real hwpx", "application/octet-stream", "sample.hwpx"
        )
        assert status in (STATUS_OK, STATUS_UNSUPPORTED, STATUS_FAILED), status
        if status == STATUS_OK:
            # 한/글이 열어준 경우 — 파싱 자체는 성공이므로 더 볼 게 없다.
            return
        # 미지원/실패면 탐지 없이 통과해야 한다(모델 없이도 확인 가능한 계약).
        assert text == ""
        assert reason
        return

    text, status, reason = parse_document(hwpx_bytes, "application/octet-stream", "sample.hwpx")
    assert status == STATUS_OK, f"실제 한/글 설치 환경에서 HWPX 파싱 실패: {reason}"
    assert "박영희" in text
    assert "park@example.com" in text

    if not model_status.all_ready():
        pytest.skip("PII/인젝션 모델이 로컬에 없어 마스킹 결과를 검증할 수 없음")
    registry.reset_cache()
    result = asyncio.run(
        run_pipeline(file_bytes=hwpx_bytes, mime_type="application/octet-stream", file_name="sample.hwpx")
    )
    assert result["scanStatus"] == STATUS_OK
    # 이메일은 실 PII 인코더가 안정적으로 잡지만, 이름은 놓치는 경우가 실측됐다
    # (모델 정확도 이슈) — 이메일만 확인한다.
    assert result["stats"]["piiCount"] >= 1
    assert "park@example.com" not in result["maskedText"]
