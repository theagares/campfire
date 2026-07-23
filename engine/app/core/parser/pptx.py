"""
app/core/parser/pptx.py
PPTX 파싱 — python-pptx (PLAN §6).

모든 슬라이드의 텍스트 프레임(+표)을 순회해 하나의 문서 텍스트로 합친다.
"""

from __future__ import annotations

import io


def extract_pptx(file_bytes: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))

    slides: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = "\n".join(
                    p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()
                )
                if text:
                    parts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        parts.append(" | ".join(cells))
        if parts:
            slides.append(f"[슬라이드 {idx}]\n" + "\n".join(parts))

    return "\n\n---\n\n".join(slides)
